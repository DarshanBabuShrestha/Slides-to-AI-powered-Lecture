from fastapi import FastAPI, File, UploadFile, Form,HTTPException
from fastapi.responses import JSONResponse
import fitz  # PyMuPDF for PDF text extraction
from pptx import Presentation
import os
import requests
from google.cloud import texttospeech
from dotenv import load_dotenv
import shutil
import uuid
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
import random
app = FastAPI()
import json

UPLOAD_DIR = "uploads"
os.makedirs(UPLOAD_DIR, exist_ok=True)
app.mount("/uploads", StaticFiles(directory=UPLOAD_DIR), name="uploads")

EXTRACTED_TEXT = ""

@app.get("/")
def root():
    return {"message": "FastAPI is running!"}

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Allow frontend access
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

client = texttospeech.TextToSpeechClient()
load_dotenv()

def extract_text_from_pdf(file_path):
    global EXTRACTED_TEXT
    doc = fitz.open(file_path)
    text = "\n".join(page.get_text() for page in doc)
    EXTRACTED_TEXT = text.strip() if text.strip() else "Error: No readable text found in PDF."
    return EXTRACTED_TEXT

def extract_text_from_pptx(file_path):
    global EXTRACTED_TEXT
    prs = Presentation(file_path)
    text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
    EXTRACTED_TEXT = text.strip() if text.strip() else "Error: No readable text found in PPTX."
    return EXTRACTED_TEXT

def generate_lecture_transcript_gemini(extracted_text):
    if not extracted_text or "Error" in extracted_text:
        return "Error: No valid text found in the document."

    gemini_api_url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-1.5-flash:generateContent?key={os.getenv('GEMINI_API_KEY')}"

    payload = {
        "contents": [{
            "parts": [{
                "text": f"""You are an AI lecture assistant delivering a **structured, engaging, and interactive lecture** based on the provided slides. Your goal is to generate a **natural, student-friendly AI narration** that feels exactly like a real professor explaining the topic in a live classroom.  

    <speak>  

    <break time="1000ms"/>  

    ### **Introduction**  

    Alright class, let's begin. <break time="700ms"/>  

    Today, we are going to explore something really interesting. <break time="500ms"/>  

    Have you ever wondered... <break time="1000ms"/> how artificial intelligence actually learns? <break time="700ms"/>  

    Let's break it down step by step.  

    ---

    ### **Slide-Based Explanation**  

    <break time="1200ms"/>  
    **Now, let's move to Slide 1.** <break time="700ms"/>  

    This slide introduces the concept of <prosody rate="85%"> Machine Learning. </prosody>  
    At its core, machine learning is about making computers learn from data.  

    Think of it like this... <break time="700ms"/>  
    When you play a video game, the more you play, the better you get. <break time="500ms"/>  

    Computers do the same thing – they improve over time, based on data.  

    ---

    <break time="1500ms"/>  
    ### **Mathematical Symbols Explained Naturally**  

    Alright, let's go to **Slide 2**, where we see an important equation:  

    <break time="1000ms"/>  
    Newton's Second Law of Motion.  

    On the slide, you see this equation:  

    <break time="700ms"/>  
    F equals m times a.  

    Notice how I said that? <break time="500ms"/>  

    Instead of saying "asterisk" or "star", I simply said **"times."**  

    So remember, when we see:  

    <break time="700ms"/>  
    **F = m * a**  

    We say:  

    <break time="700ms"/>  
    **"Force is equal to mass multiplied by acceleration."**  

    That makes it sound **natural** and **easy to understand.**  

    ---

    <break time="1500ms"/>  
    ### **Pseudocode Explained Like a Professor**  

    Now, let's go to **Slide 3**. <break time="700ms"/>  

    Here, we see a basic programming logic.  

    On the slide, you see something like this:  

    <break time="700ms"/>  
    "If X is greater than 10, print 'High'."  

    We don’t say **"If open parenthesis X greater than 10 close parenthesis..."** <break time="700ms"/>  

    That sounds **robotic and unnatural.**  

    Instead, a professor would say:  

    <break time="700ms"/>  
    "If X is greater than 10, then the system prints 'High' on the screen."  

    Sounds better, right?  

    ---

    <break time="1500ms"/>  
    ### **Real-World Example to Engage Students**  

    Alright, let's think of an example.  

    <break time="700ms"/>  

    You know how Netflix recommends shows based on what you watch?  

    That's machine learning! <break time="700ms"/>  

    It learns your preferences over time, just like how we get better at playing a game.  

    ---

    <break time="1200ms"/>  
    ### **Encouraging Student Engagement**  

    Now, if any of this isn’t clear, feel free to **ask the chatbot** for another explanation.  

    <break time="700ms"/>  

    Or, if you want a **real-world example**, just type your question in the chatbot!  

    ---

    <break time="2000ms"/>  
    ### **Conclusion**  

    Okay, before we wrap up, let's summarize:  

    <break time="700ms"/>  
    - Machine learning helps computers learn from data.  
    - **Equations should be narrated naturally** (F = m * a → "Force is mass times acceleration").  
    - **Pseudocode should be explained in a spoken way** (Not just reading code literally).  

    ---

    Alright, that's it for today.  

    <break time="1000ms"/>  

    See you in the next lecture!  

    </speak>  

    Now, generate an **engaging, student-friendly lecture transcript** based on the provided slides, while following this **natural narration format**.  

    \n\n{extracted_text}"""
            }]
        }]
    }



    headers = {"Content-Type": "application/json"}

    try:
        response = requests.post(gemini_api_url, json=payload, headers=headers, timeout=15)
        response_data = response.json()

        if response.status_code == 200 and "candidates" in response_data:
            return response_data["candidates"][0]["content"]["parts"][0]["text"].strip()
        else:
            return "Error generating transcript"

    except Exception as e:
        print(f" Gemini API Error: {e}")
        return "Error generating transcript"

def text_to_speech(transcript):
    audio_filename = f"lecture_{uuid.uuid4().hex}.mp3"
    output_audio_path = os.path.join(UPLOAD_DIR, audio_filename)

    synthesis_input = texttospeech.SynthesisInput(text=transcript)

    voice = texttospeech.VoiceSelectionParams(
        language_code="en-US",
        name="en-US-Studio-O",
    )

    audio_config = texttospeech.AudioConfig(
        audio_encoding=texttospeech.AudioEncoding.MP3,
        speaking_rate=1.0
    )

    try:
        print(" Generating AI Audio...")
        response = client.synthesize_speech(
            input=synthesis_input, voice=voice, audio_config=audio_config
        )

        print(" TTS Response received.")

        with open(output_audio_path, "wb") as out:
            out.write(response.audio_content)

        return audio_filename  

    except Exception as e:
        print(f" ERROR generating TTS audio: {e}")
        return None

@app.post("/upload")
async def upload_file(file: UploadFile = File(...)):
    global EXTRACTED_TEXT
    try:
        file_ext = file.filename.split('.')[-1].lower()
        file_path = os.path.join(UPLOAD_DIR, file.filename)

        with open(file_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)

        if not os.path.exists(file_path):
            print(f" ERROR: File {file_path} was not saved properly!")
            return JSONResponse(content={"message": "File save failed."}, status_code=500)

        print(f" File successfully saved at: {file_path}")

        extracted_text = extract_text_from_pdf(file_path) if file_ext == 'pdf' else extract_text_from_pptx(file_path)

        ai_transcript = generate_lecture_transcript_gemini(extracted_text)
        if "Error" in ai_transcript:
            return JSONResponse(content={"message": "Transcript generation failed."}, status_code=400)

        audio_filename = text_to_speech(ai_transcript)
        if not audio_filename:
            return JSONResponse(content={"message": "Audio generation failed."}, status_code=500)

        file_url = f"http://127.0.0.1:8080/uploads/{file.filename}"
        audio_url = f"http://127.0.0.1:8080/uploads/{audio_filename}"

        print(f" Returning JSON Response:\nFile URL: {file_url}\nAudio URL: {audio_url}")  # Debugging log

        return JSONResponse(content={
            "file_url": file_url,
            "audio_url": audio_url,
        }, status_code=200)

    except Exception as e:
        print(f" FastAPI Error: {e}")  
        return JSONResponse(content={"error": str(e)}, status_code=500)

@app.post("/chat")
async def chat(query: str = Form(...)):
    global EXTRACTED_TEXT

    if not EXTRACTED_TEXT:
        return JSONResponse(content={
            "response": "I don't have access to the course material. Please upload a PDF or PPTX first.",
            "audio_url": None
        })

    if not query.strip():
        return JSONResponse(content={"response": "Please ask a valid question.", "audio_url": None})

    gemini_api_url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-1.5-flash:generateContent?key={os.getenv('GEMINI_API_KEY')}"

    payload = {
        "contents": [{
            "parts": [{
                "text": f"""
                You are an AI Proferssor with access to the following document:

                ---
                {EXTRACTED_TEXT}
                ---

                Your primary task is to answer the question from anywhere when the user askes
                question about the document. The question should be stricly from the document 
                provided

                **Question:** {query}
                """
            }]
        }]
    }

    headers = {"Content-Type": "application/json"}

    try:
        response = requests.post(gemini_api_url, json=payload, headers=headers, timeout=15)
        response_data = response.json()

        if response.status_code == 200 and "candidates" in response_data:
            response_text = response_data["candidates"][0]["content"]["parts"][0]["text"].strip()
        else:
            response_text = "I couldn't find a relevant answer in the course material."

    except Exception as e:
        print(f" Error calling Gemini API: {e}")
        response_text = "Error generating AI response."

    audio_filename = text_to_speech(response_text) if response_text else None
    audio_url = f"http://127.0.0.1:8080/uploads/{audio_filename}" if audio_filename else None

    return JSONResponse(content={"response": response_text, "audio_url": audio_url})

@app.post("/generate_quiz")
async def generate_quiz():
    global EXTRACTED_TEXT

    if not EXTRACTED_TEXT:
        raise HTTPException(status_code=400, detail="No transcript found. Upload slides first.")

    gemini_api_url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-1.5-flash:generateContent?key={os.getenv('GEMINI_API_KEY')}"

    payload = {
        "contents": [{
            "parts": [{
                "text": f"Generate a 10-question quiz from the following transcript. Ensure a mix of multiple-choice , true/false questions or anything that can fit under the criteria of multiple choice question. Format it as a **valid JSON array**, where each object contains 'question', 'options', and 'answer'. Do **not** return markdown formatting (```json) or any code blocks:\n\n{EXTRACTED_TEXT}"
            }]
        }]
    }

    headers = {"Content-Type": "application/json"}

    try:
        response = requests.post(gemini_api_url, json=payload, headers=headers, timeout=15)
        response_data = response.json()

        if response.status_code == 200 and "candidates" in response_data:
            raw_quiz_text = response_data["candidates"][0]["content"]["parts"][0]["text"].strip()

            try:
                quiz_data = json.loads(raw_quiz_text)
            except json.JSONDecodeError:
                return JSONResponse(content={"message": "Invalid AI response format"}, status_code=500)

            return JSONResponse(content={"quiz": quiz_data})
        else:
            return JSONResponse(content={"message": "Quiz generation failed."}, status_code=400)

    except Exception as e:
        print(f" Error generating quiz: {e}")
        return JSONResponse(content={"error": "Quiz generation error."}, status_code=500)
