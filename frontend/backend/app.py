from fastapi import FastAPI, File, UploadFile
from fastapi.responses import JSONResponse
import fitz  # PyMuPDF for PDF text extraction
from pptx import Presentation
import os
import requests
from google.cloud import texttospeech
from dotenv import load_dotenv
import shutil
import uuid
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles

app = FastAPI()

UPLOAD_DIR = "uploads"
os.makedirs(UPLOAD_DIR, exist_ok=True)
app.mount("/uploads", StaticFiles(directory=UPLOAD_DIR), name="uploads")
@app.get("/")
def root():
    return {"message": "FastAPI is running!"}
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000"],  # Allow frontend access
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

client = texttospeech.TextToSpeechClient()

load_dotenv()

def extract_text_from_pdf(file_path):
    doc = fitz.open(file_path)
    text = "\n".join(page.get_text() for page in doc)
    return text.strip() if text.strip() else "Error: No readable text found in PDF."

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
    return text.strip() if text.strip() else "Error: No readable text found in PPTX."

def generate_lecture_transcript_gemini(extracted_text):
    if not extracted_text or "Error" in extracted_text:
        return "Error: No valid text found in the document."

    gemini_api_url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-1.5-flash:generateContent?key={os.getenv('GEMINI_API_KEY')}"

    payload = {
        "contents": [{
            "parts": [{
                "text": f"Convert the following slide text into a structured, well-explained lecture transcript that feels like a live lecture. Ensure it is natural-sounding and suitable for AI voice narration and make it below 4800 bytes:\n\n{extracted_text}"
            }]
        }]
    }

    headers = {"Content-Type": "application/json"}

    try:
        response = requests.post(gemini_api_url, json=payload, headers=headers, timeout=15)
        response_data = response.json()

        if response.status_code == 200 and "candidates" in response_data:
            return response_data["candidates"][0]["content"]["parts"][0]["text"].strip()
        else:
            return "Error generating transcript"

    except Exception as e:
        print(f"❌ Gemini API Error: {e}")
        return "Error generating transcript"

def text_to_speech(transcript):
    audio_filename = f"lecture_{uuid.uuid4().hex}.mp3"
    output_audio_path = os.path.join(UPLOAD_DIR, audio_filename)

    synthesis_input = texttospeech.SynthesisInput(text=transcript)

    voice = texttospeech.VoiceSelectionParams(
        language_code="en-US",
        name="en-US-Studio-O",
    )

    audio_config = texttospeech.AudioConfig(
        audio_encoding=texttospeech.AudioEncoding.MP3,
        speaking_rate=1.0
    )

    try:
        print("🚀 Generating AI Audio...")
        response = client.synthesize_speech(
            input=synthesis_input, voice=voice, audio_config=audio_config
        )

        print("✅ TTS Response received.")

        with open(output_audio_path, "wb") as out:
            out.write(response.audio_content)

        return audio_filename  

    except Exception as e:
        print(f"❌ ERROR generating TTS audio: {e}")
        return None

@app.post("/upload")
async def upload_file(file: UploadFile = File(...)):
    try:
        file_ext = file.filename.split('.')[-1].lower()
        file_path = os.path.join(UPLOAD_DIR, file.filename)

        with open(file_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)

        if not os.path.exists(file_path):
            print(f"❌ ERROR: File {file_path} was not saved properly!")
            return JSONResponse(content={"message": "File save failed."}, status_code=500)

        print(f"✅ File successfully saved at: {file_path}")

        # ✅ Extract text from PDF or PPTX
        extracted_text = extract_text_from_pdf(file_path) if file_ext == 'pdf' else extract_text_from_pptx(file_path)

        # ✅ Generate AI transcript
        ai_transcript = generate_lecture_transcript_gemini(extracted_text)
        if "Error" in ai_transcript:
            return JSONResponse(content={"message": "Transcript generation failed."}, status_code=400)

        # ✅ Generate AI audio
        audio_filename = text_to_speech(ai_transcript)
        if not audio_filename:
            return JSONResponse(content={"message": "Audio generation failed."}, status_code=500)

        # ✅ Ensure the correct `uploads` URL
        file_url = f"http://127.0.0.1:8080/uploads/{file.filename}"
        audio_url = f"http://127.0.0.1:8080/uploads/{audio_filename}"

        print(f"✅ Returning JSON Response:\nFile URL: {file_url}\nAudio URL: {audio_url}")  # Debugging log

        return JSONResponse(content={
            "file_url": file_url,
            "audio_url": audio_url,
            
        }, status_code=200)

    except Exception as e:
        print(f"❌ FastAPI Error: {e}")  
        return JSONResponse(content={"error": str(e)}, status_code=500)

@app.get("/")
def read_root():
    return {"message": "Hello World"}

@app.get("/test-tts")
def test_tts():
    test_transcript = "This is a test of Google Cloud Text-to-Speech inside FastAPI."
    audio_filename = text_to_speech(test_transcript)

    if audio_filename:
        return {"message": "TTS worked!", "audio_url": f"http://127.0.0.1:8080/uploads/{audio_filename}"}
    else:
        return {"message": "TTS failed!"}